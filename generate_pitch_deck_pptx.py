#!/usr/bin/env python3
"""
Vidya Nirmaan — Editable Pitch Deck (.pptx) Generator
Same content as the PDF pitch deck but in editable PowerPoint format.
"""

import os
import tempfile
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, Circle
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE

# ─── Colors ───
PRIMARY = RGBColor(0x15, 0x65, 0xC0)
PRIMARY_DARK = RGBColor(0x0D, 0x47, 0xA1)
PRIMARY_LIGHT = RGBColor(0xE3, 0xF2, 0xFD)
ORANGE = RGBColor(0xF5, 0x7C, 0x00)
PURPLE = RGBColor(0x7B, 0x1F, 0xA2)
RED = RGBColor(0xD3, 0x2F, 0x2F)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
TEAL = RGBColor(0x00, 0x89, 0x7B)
AMBER = RGBColor(0xF9, 0xA8, 0x25)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x75, 0x75, 0x75)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
BROWN = RGBColor(0x5D, 0x40, 0x37)

TEMP_DIR = tempfile.mkdtemp()

def hex_to_rgb(h):
    h = h.lstrip('#')
    return tuple(int(h[i:i+2], 16)/255.0 for i in (0, 2, 4))


# ─── Helper Functions ───

def add_textbox(slide, left, top, width, height, text, font_size=12, bold=False,
                color=DARK, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a simple text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_shape_with_text(slide, left, top, width, height, text, fill_color, text_color=WHITE,
                        font_size=12, bold=True, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, subtitle=None):
    """Add a colored rounded rectangle with centered text."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    p.font.name = 'Calibri'

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(font_size - 3)
        p2.font.color.rgb = text_color
        p2.font.name = 'Calibri'
        p2.alignment = PP_ALIGN.CENTER

    # Center vertically
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)

    return shape


def add_table(slide, left, top, width, height, data, col_widths=None,
              header_color=PRIMARY, header_text_color=WHITE):
    """Add a styled table to a slide."""
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.name = 'Calibri'

                if row_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = header_text_color
                    paragraph.alignment = PP_ALIGN.CENTER
                elif col_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = PRIMARY
                else:
                    paragraph.font.color.rgb = DARK

            # Header row background
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


def add_footer(slide, page_num, total):
    """Add footer bar to slide."""
    # Dark blue bar at bottom
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0), Inches(7.15), Inches(13.33), Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_DARK
    shape.line.fill.background()

    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = f'Koreidea Innovations  |  Vidya Nirmaan  |  IndiaAI Innovation Challenge  |  {page_num} / {total}'
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER


def add_section_label(slide, text, top=Inches(0.4)):
    """Add orange section label at top."""
    add_textbox(slide, Inches(0.6), top, Inches(5), Inches(0.3),
                text, font_size=11, bold=True, color=ORANGE)


def add_slide_title(slide, text, top=Inches(0.65)):
    """Add large blue slide title."""
    add_textbox(slide, Inches(0.6), top, Inches(12), Inches(0.6),
                text, font_size=28, bold=True, color=PRIMARY_DARK)


def add_slide_subtitle(slide, text, top=Inches(1.25)):
    """Add blue subtitle text."""
    add_textbox(slide, Inches(0.6), top, Inches(11.5), Inches(0.6),
                text, font_size=13, color=PRIMARY)


# ─── Diagram Generators (reused from PDF script) ───

def generate_problem_visual():
    fig, ax = plt.subplots(figsize=(11, 4.5))
    ax.set_xlim(0, 11)
    ax.set_ylim(0, 4.5)
    ax.axis('off')
    fig.patch.set_facecolor('white')

    problems = [
        ('Manual\nPaper-Based', 'Demand forms processed\nmanually across 707 mandals', '#D32F2F', 0.3),
        ('No Fraud\nDetection', 'Demand inflation goes\nundetected, wasting crores', '#E65100', 3.0),
        ('Inequitable\nAllocation', 'No objective prioritization\nof 45,000+ schools', '#F57C00', 5.7),
        ('No Future\nPlanning', 'No enrolment forecasting\nfor proactive infra planning', '#F9A825', 8.4),
    ]

    for label, desc, color, x in problems:
        circle = Circle((x+1.1, 3.2), 0.65, facecolor=hex_to_rgb(color), edgecolor='white', linewidth=2, zorder=3)
        ax.add_patch(circle)
        ax.text(x+1.1, 3.2, label, ha='center', va='center', fontsize=9, fontweight='bold', color='white', zorder=4)
        box = FancyBboxPatch((x+0.05, 1.2), 2.1, 1.2, boxstyle="round,pad=0.08",
                              facecolor='white', edgecolor=hex_to_rgb(color), linewidth=1.5)
        ax.add_patch(box)
        ax.text(x+1.1, 1.8, desc, ha='center', va='center', fontsize=8, color='#333', linespacing=1.4)
        ax.plot([x+1.1, x+1.1], [2.55, 2.4], color=hex_to_rgb(color), linewidth=2)

    bar = FancyBboxPatch((0.5, 0.15), 10, 0.7, boxstyle="round,pad=0.05",
                          facecolor=hex_to_rgb('#0D47A1'), edgecolor='none', alpha=0.9)
    ax.add_patch(bar)
    ax.text(5.5, 0.5, '45,000+ Schools  |  57 Districts  |  707 Mandals  |  Andhra Pradesh  |  Dept. of School Education',
            ha='center', va='center', fontsize=10, fontweight='bold', color='white')

    path = os.path.join(TEMP_DIR, 'problem_visual.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path


def generate_solution_overview():
    fig, ax = plt.subplots(figsize=(11, 5.5))
    ax.set_xlim(-0.5, 12.0)
    ax.set_ylim(0, 5.5)
    ax.axis('off')
    fig.patch.set_facecolor('white')

    center = FancyBboxPatch((4.05, 2.0), 3.4, 1.6, boxstyle="round,pad=0.1",
                             facecolor=hex_to_rgb('#1565C0'), edgecolor=hex_to_rgb('#0D47A1'), linewidth=2)
    ax.add_patch(center)
    ax.text(5.75, 3.15, 'Vidya Nirmaan', ha='center', va='center', fontsize=15, fontweight='bold', color='white')
    ax.text(5.75, 2.7, 'AI-Powered Mobile App', ha='center', va='center', fontsize=10, color='#B3D4FC')
    ax.text(5.75, 2.3, 'Flutter + Supabase + FastAPI', ha='center', va='center', fontsize=9, color='#B3D4FC')

    features = [
        ('AI Priority\nScoring', '4-Factor Weighted\n0-100 Composite', '#7B1FA2', 1.2, 4.5),
        ('Fraud\nDetection', 'Isolation Forest\n7 Validation Rules', '#D32F2F', 4.4, 4.5),
        ('Enrolment\nForecasting', 'Linear Regression\nCohort Progression', '#F57C00', 7.7, 4.5),
        ('3-Stage\nApproval', 'AI > Assessment\n> Officer', '#00897B', 9.5, 2.8),
        ('Offline\nFirst', '5 Hive Boxes\nAuto-Sync', '#4CAF50', 7.7, 0.5),
        ('5 User\nRoles', 'State to School\nRBAC Enforced', '#1565C0', 4.4, 0.3),
        ('Bilingual', 'English + Telugu\nFull Localization', '#F9A825', 1.2, 0.5),
        ('Field\nInspection', '50+ Data Points\nGPS Verified', '#5D4037', 0.2, 2.8),
    ]

    for label, desc, color, x, y in features:
        box = FancyBboxPatch((x-0.1, y-0.1), 2.2, 1.0, boxstyle="round,pad=0.05",
                              facecolor='white', edgecolor=hex_to_rgb(color), linewidth=1.5)
        ax.add_patch(box)
        ax.text(x+1.0, y+0.6, label, ha='center', va='center', fontsize=9, fontweight='bold', color=hex_to_rgb(color))
        ax.text(x+1.0, y+0.2, desc, ha='center', va='center', fontsize=7, color='#555', linespacing=1.3)
        cx, cy = 5.75, 2.8
        fx, fy = x+1.0, y+0.45
        ax.plot([fx, cx], [fy, cy], color=hex_to_rgb(color), linewidth=0.8, alpha=0.35, linestyle='--')

    path = os.path.join(TEMP_DIR, 'solution_overview.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path


def generate_ai_models():
    fig, ax = plt.subplots(figsize=(11, 5))
    ax.set_xlim(0, 11)
    ax.set_ylim(0, 5)
    ax.axis('off')
    fig.patch.set_facecolor('white')

    models = [
        {'title': 'Priority Scoring', 'sub': 'School Prioritization', 'color': '#7B1FA2', 'x': 0.2,
         'items': ['Enrolment Pressure (30%)', 'Infrastructure Gap (30%)', 'CWSN Needs (20%)', 'Accessibility (20%)'],
         'output': 'Score 0-100\nCRITICAL > HIGH > MEDIUM > LOW'},
        {'title': 'Demand Validation', 'sub': 'Fraud Detection', 'color': '#D32F2F', 'x': 3.8,
         'items': ['7 Rule-Based Checks', 'Isolation Forest ML', 'Unit Cost Deviation', 'Peer Comparison'],
         'output': 'APPROVED (>=80)\nFLAGGED (50-79) | REJECTED (<50)'},
        {'title': 'Enrolment Forecast', 'sub': '3-Year Prediction', 'color': '#F57C00', 'x': 7.4,
         'items': ['Linear Regression', 'Cohort Progression', 'Confidence Decay', 'Backend + Client Fallback'],
         'output': 'Y1: ~93% conf\nY2: ~85% | Y3: ~76%'},
    ]

    for m in models:
        x, color = m['x'], m['color']
        hdr = FancyBboxPatch((x, 3.5), 3.2, 1.1, boxstyle="round,pad=0.05",
                              facecolor=hex_to_rgb(color), edgecolor=hex_to_rgb(color), linewidth=1.5)
        ax.add_patch(hdr)
        ax.text(x+1.6, 4.25, m['title'], ha='center', fontsize=12, fontweight='bold', color='white')
        ax.text(x+1.6, 3.85, m['sub'], ha='center', fontsize=9, color='white', alpha=0.9)

        body = FancyBboxPatch((x, 1.3), 3.2, 2.2, boxstyle="round,pad=0.05",
                               facecolor='white', edgecolor=hex_to_rgb(color), linewidth=1.5)
        ax.add_patch(body)
        for i, item in enumerate(m['items']):
            ax.text(x+0.25, 3.15 - i*0.4, f'\u2022  {item}', fontsize=9, color='#333')

        out = FancyBboxPatch((x, 0.2), 3.2, 0.9, boxstyle="round,pad=0.05",
                              facecolor=hex_to_rgb(color), edgecolor=hex_to_rgb(color), linewidth=1, alpha=0.15)
        ax.add_patch(out)
        lines = m['output'].split('\n')
        for i, line in enumerate(lines):
            ax.text(x+1.6, 0.75 - i*0.28, line, ha='center', fontsize=8.5, fontweight='bold', color=hex_to_rgb(color))

    path = os.path.join(TEMP_DIR, 'ai_models.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path


def generate_pipeline_visual():
    fig, ax = plt.subplots(figsize=(11, 3.5))
    ax.set_xlim(0, 11)
    ax.set_ylim(0, 3.5)
    ax.axis('off')
    fig.patch.set_facecolor('white')

    stages = [
        ('School HM\nRaises Demand', '#1565C0', 0.2),
        ('Stage 1\nAI Validation', '#7B1FA2', 2.9),
        ('Stage 2\nField Assessment', '#00897B', 5.6),
        ('Stage 3\nOfficer Decision', '#F57C00', 8.3),
    ]

    for label, color, x in stages:
        box = FancyBboxPatch((x, 1.2), 2.2, 1.5, boxstyle="round,pad=0.08",
                              facecolor=hex_to_rgb(color), edgecolor='white', linewidth=2)
        ax.add_patch(box)
        ax.text(x+1.1, 1.95, label, ha='center', va='center', fontsize=11, fontweight='bold', color='white')

    for i in range(3):
        x = stages[i][2] + 2.2
        ax.annotate('', xy=(x+0.7, 1.95), xytext=(x, 1.95),
                    arrowprops=dict(arrowstyle='->', lw=2.5, color='#333'))

    sublabels = [
        ('infra_type, qty\namount, year', 1.3, 0.65),
        ('7 Rules + ML\nScore 0-100', 4.0, 0.65),
        ('50+ fields, GPS\nPhoto evidence', 6.7, 0.65),
        ('APPROVE / FLAG\n/ REJECT', 9.4, 0.65),
    ]
    for text, x, y in sublabels:
        ax.text(x, y, text, ha='center', va='center', fontsize=8, color='#555', linespacing=1.3)

    gate = FancyBboxPatch((4.8, 2.9), 5.8, 0.4, boxstyle="round,pad=0.03",
                           facecolor=hex_to_rgb('#D32F2F'), edgecolor='none', alpha=0.9)
    ax.add_patch(gate)
    ax.text(7.7, 3.1, 'HARD GATE: Field assessment required before officer approval',
            ha='center', fontsize=8, fontweight='bold', color='white')

    path = os.path.join(TEMP_DIR, 'pipeline_visual.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path


def generate_impact_chart():
    fig, axes = plt.subplots(1, 2, figsize=(11, 4.5))
    fig.patch.set_facecolor('white')

    ax1 = axes[0]
    categories = ['Demand\nProcessing', 'Fraud\nDetection', 'School\nPrioritization', 'Forecast\nAccuracy']
    current = [20, 5, 15, 0]
    vidya = [95, 85, 100, 76]
    x = np.arange(len(categories))
    width = 0.35
    bars1 = ax1.bar(x - width/2, current, width, label='Current (Manual)', color=hex_to_rgb('#D32F2F'), alpha=0.7)
    bars2 = ax1.bar(x + width/2, vidya, width, label='Vidya Nirmaan (AI)', color=hex_to_rgb('#1565C0'), alpha=0.9)
    ax1.set_ylabel('Efficiency %', fontsize=11, fontweight='bold')
    ax1.set_title('Manual vs AI-Powered', fontsize=13, fontweight='bold', color=hex_to_rgb('#0D47A1'))
    ax1.set_xticks(x)
    ax1.set_xticklabels(categories, fontsize=9)
    ax1.set_ylim(0, 120)
    ax1.legend(fontsize=8, loc='upper left', framealpha=0.9)
    ax1.spines['top'].set_visible(False)
    ax1.spines['right'].set_visible(False)
    for bar in bars1:
        h = bar.get_height()
        ax1.text(bar.get_x() + bar.get_width()/2., h + 2, f'{int(h)}%', ha='center', fontsize=8, color='#555')
    for bar in bars2:
        h = bar.get_height()
        ax1.text(bar.get_x() + bar.get_width()/2., h + 2, f'{int(h)}%', ha='center', fontsize=8, fontweight='bold', color=hex_to_rgb('#0D47A1'))

    ax2 = axes[1]
    phases = ['Pilot\n(Current)', 'Phase 1\n(3 Months)', 'Phase 2\n(6 Months)']
    schools = [319, 15000, 45000]
    colors_bar = [hex_to_rgb('#1565C0'), hex_to_rgb('#F57C00'), hex_to_rgb('#4CAF50')]
    bars = ax2.bar(phases, schools, color=colors_bar, alpha=0.85, edgecolor='white', linewidth=1.5)
    ax2.set_ylabel('Schools Covered', fontsize=11, fontweight='bold')
    ax2.set_title('Scalability Roadmap', fontsize=13, fontweight='bold', color=hex_to_rgb('#0D47A1'))
    ax2.spines['top'].set_visible(False)
    ax2.spines['right'].set_visible(False)
    for bar, val in zip(bars, schools):
        h = bar.get_height()
        ax2.text(bar.get_x() + bar.get_width()/2., h + 800, f'{val:,}', ha='center', fontsize=10, fontweight='bold', color=hex_to_rgb('#0D47A1'))

    plt.tight_layout(pad=2)
    path = os.path.join(TEMP_DIR, 'impact_chart.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path


def generate_tech_stack_visual():
    fig, ax = plt.subplots(figsize=(12, 4))
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 4)
    ax.axis('off')
    fig.patch.set_facecolor('white')

    layers = [
        ('Mobile App', 'Flutter 3.10+  |  Riverpod 3.x  |  fl_chart  |  flutter_map  |  Material Design 3', '#1565C0', 3.0),
        ('State Mgmt', 'Riverpod Notifier/AsyncValue  |  Role-Based Provider Scoping  |  Offline Fallback', '#1976D2', 2.2),
        ('Backend', 'Supabase (PostgreSQL + RLS + Auth)  |  FastAPI ML (Railway)  |  scikit-learn', '#4CAF50', 1.4),
        ('Offline Layer', 'Hive (5 Boxes)  |  Auto-Sync  |  Queue & Retry  |  Graceful Degradation', '#00897B', 0.6),
    ]

    for label, desc, color, y in layers:
        box = FancyBboxPatch((0.3, y), 11.4, 0.65, boxstyle="round,pad=0.05",
                              facecolor=hex_to_rgb(color), edgecolor='white', linewidth=2, alpha=0.9)
        ax.add_patch(box)
        ax.text(1.7, y+0.33, label, va='center', fontsize=11, fontweight='bold', color='white')
        ax.text(4.2, y+0.33, desc, va='center', fontsize=9, color='white', alpha=0.95)

    path = os.path.join(TEMP_DIR, 'tech_stack.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path


# ─── PPTX Builder ───

def build_pptx(diagrams):
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                'Vidya_Nirmaan_Pitch_Deck.pptx')

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Blank layout
    TOTAL_SLIDES = 12

    # ═══════════════════════════════════════
    # SLIDE 1: COVER
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)

    # Title
    add_textbox(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(1),
                'Vidya Nirmaan', font_size=44, bold=True, color=PRIMARY_DARK, alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, Inches(1.5), Inches(2.1), Inches(10), Inches(0.6),
                'AI-Powered School Infrastructure Planning & Monitoring', font_size=20,
                color=PRIMARY, alignment=PP_ALIGN.CENTER)

    # Orange line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(2.85), Inches(4), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.fill.background()

    # Badge
    add_textbox(slide, Inches(1.5), Inches(3.1), Inches(10), Inches(0.4),
                'IndiaAI Innovation Challenge \u2014 Problem Statement 5', font_size=15,
                bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.3),
                'Department of School Education, Andhra Pradesh', font_size=13,
                color=GRAY, alignment=PP_ALIGN.CENTER)

    # Metrics row
    metrics = [('319', 'Schools'), ('57', 'Districts'), ('707', 'Mandals'), ('3', 'AI Models'), ('5', 'User Roles')]
    for i, (num, label) in enumerate(metrics):
        x = Inches(1.8 + i * 2.0)
        add_textbox(slide, x, Inches(4.3), Inches(1.5), Inches(0.6),
                    num, font_size=28, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(4.85), Inches(1.5), Inches(0.3),
                    label, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Bottom info
    add_textbox(slide, Inches(1.5), Inches(5.6), Inches(10), Inches(0.3),
                'March 2026  |  Flutter + Supabase + FastAPI  |  Offline-First Mobile App',
                font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_footer(slide, 1, TOTAL_SLIDES)

    # ═══════════════════════════════════════
    # SLIDE 2: THE PROBLEM
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_section_label(slide, 'THE PROBLEM')
    add_slide_title(slide, 'Reimagining School Infrastructure Planning')
    add_slide_subtitle(slide, 'Andhra Pradesh manages 45,000+ government schools with manual, paper-based infrastructure planning. No automated fraud detection, no objective prioritization, no future-readiness.')
    slide.shapes.add_picture(diagrams['problem'], Inches(0.8), Inches(2.2), Inches(11.5), Inches(4.5))
    add_footer(slide, 2, TOTAL_SLIDES)

    # ═══════════════════════════════════════
    # SLIDE 3: THE SOLUTION
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_section_label(slide, 'OUR SOLUTION')
    add_slide_title(slide, 'Vidya Nirmaan \u2014 End-to-End AI Platform')
    add_slide_subtitle(slide, 'A mobile-first application that digitizes the entire infrastructure lifecycle: demand \u2192 validation \u2192 assessment \u2192 approval, powered by 3 AI/ML models.')
    slide.shapes.add_picture(diagrams['solution'], Inches(0.6), Inches(2.0), Inches(12), Inches(5.0))
    add_footer(slide, 3, TOTAL_SLIDES)

    # ═══════════════════════════════════════
    # SLIDE 4: HOW IT WORKS
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_section_label(slide, 'HOW IT WORKS')
    add_slide_title(slide, '3-Stage Demand Approval Pipeline')
    add_slide_subtitle(slide, 'Every infrastructure demand passes through AI validation, field verification, and officer approval. A hard gate ensures no demand is approved without ground-truth inspection.')
    slide.shapes.add_picture(diagrams['pipeline'], Inches(0.5), Inches(2.2), Inches(12), Inches(3.3))

    # Differentiators row
    diffs = [
        ('AI-First Validation', '7 automated rules +\nIsolation Forest ML'),
        ('Ground-Truth Verification', '50+ inspection fields\nwith GPS verification'),
        ('Officer Oversight', 'Manual approve/flag/reject\nwith justification notes'),
        ('Full Audit Trail', 'Every stage timestamped\nand tracked'),
    ]
    for i, (title, desc) in enumerate(diffs):
        x = Inches(0.8 + i * 3.1)
        add_textbox(slide, x, Inches(5.65), Inches(2.8), Inches(0.3),
                    title, font_size=11, bold=True, color=PRIMARY_DARK, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(5.95), Inches(2.8), Inches(0.5),
                    desc, font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_footer(slide, 4, TOTAL_SLIDES)

    # ═══════════════════════════════════════
    # SLIDE 5: AI/ML MODELS
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_section_label(slide, 'AI / ML INNOVATION')
    add_slide_title(slide, 'Three AI Models Working in Concert')
    add_slide_subtitle(slide, 'Client-side + backend ML: priority scoring for objective allocation, Isolation Forest for fraud detection, dual-model forecasting for future planning.')
    slide.shapes.add_picture(diagrams['ai_models'], Inches(0.5), Inches(2.1), Inches(12), Inches(4.8))
    add_footer(slide, 5, TOTAL_SLIDES)

    # ═══════════════════════════════════════
    # SLIDE 6: TECHNOLOGY
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_section_label(slide, 'TECHNOLOGY')
    add_slide_title(slide, 'Modern, Scalable Architecture')
    slide.shapes.add_picture(diagrams['tech_stack'], Inches(0.5), Inches(1.5), Inches(12), Inches(3.5))

    techs = [
        ('Cross-Platform', 'Single codebase for\niOS + Android', 'Flutter + Dart'),
        ('Real-Time Backend', 'Supabase PostgreSQL\nwith Row Level Security', 'Supabase + RLS'),
        ('ML Pipeline', 'scikit-learn models on\nRailway cloud', 'FastAPI + Python'),
        ('Offline-First', 'Hive local storage\nwith auto-sync', 'Hive + Queue'),
    ]
    for i, (title, desc, tech) in enumerate(techs):
        x = Inches(0.8 + i * 3.1)
        add_textbox(slide, x, Inches(5.2), Inches(2.8), Inches(0.3),
                    title, font_size=12, bold=True, color=PRIMARY_DARK, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(5.5), Inches(2.8), Inches(0.5),
                    desc, font_size=9, color=DARK, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(6.05), Inches(2.8), Inches(0.3),
                    tech, font_size=8, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_footer(slide, 6, TOTAL_SLIDES)

    # ═══════════════════════════════════════
    # SLIDE 7: KEY METRICS
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_section_label(slide, 'BY THE NUMBERS')
    add_slide_title(slide, 'Platform at a Glance')

    all_metrics = [
        [('319', 'Schools in\nPilot Dataset', PRIMARY), ('4,638', 'Enrolment Records\nAnalyzed', PRIMARY),
         ('799', 'Demand Plans\nValidated', PRIMARY), ('50+', 'Inspection\nData Points', PRIMARY)],
        [('7', 'Validation Rules\n+ Isolation Forest', ORANGE), ('3', 'AI/ML Models\nDeployed', ORANGE),
         ('5', 'User Roles with\nRBAC Enforcement', ORANGE), ('2', 'Languages\nEN + Telugu', ORANGE)],
        [('100%', 'Offline Capable\nFull Functionality', TEAL), ('< 2s', 'AI Validation\nResponse Time', TEAL),
         ('0', 'Manual Steps in\nFraud Detection', TEAL), ('45,000+', 'Schools\nScalability Target', TEAL)],
    ]

    for row_idx, row in enumerate(all_metrics):
        for col_idx, (num, label, color) in enumerate(row):
            x = Inches(1.0 + col_idx * 3.0)
            y_base = Inches(1.6 + row_idx * 1.8)
            add_textbox(slide, x, y_base, Inches(2.5), Inches(0.7),
                        num, font_size=34, bold=True, color=color, alignment=PP_ALIGN.CENTER)
            add_textbox(slide, x, y_base + Inches(0.6), Inches(2.5), Inches(0.5),
                        label, font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_footer(slide, 7, TOTAL_SLIDES)

    # ═══════════════════════════════════════
    # SLIDE 8: IMPACT & SCALABILITY
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_section_label(slide, 'IMPACT & SCALABILITY')
    add_slide_title(slide, 'Transforming Infrastructure Planning')
    slide.shapes.add_picture(diagrams['impact'], Inches(0.5), Inches(1.5), Inches(12), Inches(4.5))

    impacts = [
        ('Reduced Processing Time', 'From weeks to seconds\nwith AI validation'),
        ('Automated Fraud Detection', 'Isolation Forest catches\nanomalies humans miss'),
        ('Objective Prioritization', '4-factor scoring replaces\nsubjective decisions'),
        ('Proactive Planning', '3-year forecasts enable\nproactive investment'),
    ]
    for i, (title, desc) in enumerate(impacts):
        x = Inches(0.8 + i * 3.1)
        add_textbox(slide, x, Inches(6.0), Inches(2.8), Inches(0.3),
                    title, font_size=10, bold=True, color=PRIMARY_DARK, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(6.3), Inches(2.8), Inches(0.4),
                    desc, font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_footer(slide, 8, TOTAL_SLIDES)

    # ═══════════════════════════════════════
    # SLIDE 9: RESPONSIBLE AI
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_section_label(slide, 'RESPONSIBLE AI')
    add_slide_title(slide, 'Ethical, Transparent, Accountable')

    rai_data = [
        ['Principle', 'Implementation'],
        ['Transparency', 'All AI decisions show detailed rule-by-rule breakdown with scores. Officers see exactly why a demand was approved, flagged, or rejected.'],
        ['Human Oversight', '3-stage pipeline ensures human officers make final decisions. AI assists but never auto-approves without assessment + officer review.'],
        ['Fairness', 'Objective 4-factor scoring algorithm treats all schools equally. No geographic or administrative bias in prioritization.'],
        ['Privacy', 'No personal student data collected. Only aggregate enrolment counts. Supabase Row Level Security enforces data access boundaries.'],
        ['Accountability', 'Full audit trail: every validation, assessment, and officer decision is timestamped with the officer name and justification notes.'],
        ['Inclusivity', 'CWSN (Children With Special Needs) weighted at 20% in priority scoring. Bilingual interface (English + Telugu) for accessibility.'],
        ['Offline Equity', 'Full offline-first design ensures rural schools with poor connectivity are not disadvantaged. Auto-sync when connected.'],
    ]
    add_table(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(5.2), rai_data,
              col_widths=[Inches(1.8), Inches(10.2)], header_color=PRIMARY)
    add_footer(slide, 9, TOTAL_SLIDES)

    # ═══════════════════════════════════════
    # SLIDE 10: DATA GOVERNANCE
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_section_label(slide, 'DATA GOVERNANCE')
    add_slide_title(slide, 'Secure, Compliant, Role-Scoped')

    dg_data = [
        ['Aspect', 'Approach', 'Details'],
        ['Data Source', 'Government Official Data', 'UDISE+ school database, Samagra Shiksha norms, AP Dept. of Education records'],
        ['Access Control', '5-Role RBAC', 'State Official sees all; District Officer sees district; Block Officer/Inspector sees mandal; HM sees own school only'],
        ['Database Security', 'Row Level Security', 'Supabase PostgreSQL RLS policies enforce data isolation per user role at the database level'],
        ['Encryption', 'In-Transit + At-Rest', 'HTTPS for all API calls. Supabase provides AES-256 encryption at rest. JWT auth tokens for sessions'],
        ['Data Minimization', 'Aggregate Only', 'No personal student data (names, IDs). Only aggregate enrolment counts (boys, girls, total per grade)'],
        ['Offline Security', 'Local-Only Cache', 'Hive cache stored on device only. No cross-device data sharing. Cache cleared on logout'],
        ['Audit Logging', 'Full Traceability', 'Every validation, assessment, and officer decision recorded with timestamp, officer name, and justification'],
        ['Consent', 'Govt Authorization', 'Data usage authorized under IndiaAI Innovation Challenge. No third-party data sharing'],
    ]
    add_table(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(5.2), dg_data,
              col_widths=[Inches(1.6), Inches(2.2), Inches(8.2)], header_color=TEAL)
    add_footer(slide, 10, TOTAL_SLIDES)

    # ═══════════════════════════════════════
    # SLIDE 11: TEAM
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_section_label(slide, 'THE TEAM')
    add_slide_title(slide, 'Built by Engineers Who Ship')

    team_data = [
        ['Role', 'Focus Area', 'Expertise'],
        ['CEO', 'Product Vision & Architecture', 'End-to-end product development, Mech Engg background, AI/ML, Mobile & Cloud'],
        ['COO', 'Operations & Strategy', 'Business operations, government project coordination, strategic planning'],
        ['Lead Engineer', 'Full-Stack Development', 'Full-stack engineering, 5+ years experience, frontend & backend'],
        ['Product Dev Engineer', 'Product Development', 'Product development engineering, prototyping, testing & deployment pipelines'],
    ]
    add_table(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(3.5), team_data,
              col_widths=[Inches(1.8), Inches(3.2), Inches(7.0)], header_color=PRIMARY_DARK)

    add_textbox(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.8),
                'Organization: EdTech startup focused on AI-powered solutions for government education \u2014 '
                'school infrastructure (Vidya Nirmaan), early childhood development (ECD monitoring), '
                'and personalized learning companions for children.',
                font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_footer(slide, 11, TOTAL_SLIDES)

    # ═══════════════════════════════════════
    # SLIDE 12: CLOSING
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)

    add_textbox(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(1),
                'Vidya Nirmaan', font_size=44, bold=True, color=PRIMARY_DARK, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(0.8),
                'Transforming school infrastructure planning from manual paper-based workflows\n'
                'to an AI-powered, transparent, accountable, and equitable system.',
                font_size=15, color=DARK, alignment=PP_ALIGN.CENTER)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(3.2), Inches(3), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.fill.background()

    closing = [('3 AI Models', 'Priority Scoring\nValidation\nForecasting'),
               ('7 Validation Rules', 'Rule-Based +\nIsolation Forest ML'),
               ('5 User Roles', 'State to School\nRole-Based Access'),
               ('Offline-First', 'Works Without\nInternet')]
    for i, (title, desc) in enumerate(closing):
        x = Inches(1.2 + i * 2.8)
        add_textbox(slide, x, Inches(3.5), Inches(2.5), Inches(0.4),
                    title, font_size=14, bold=True, color=PRIMARY_DARK, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(3.9), Inches(2.5), Inches(0.7),
                    desc, font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.4),
                'Ready for Pilot Deployment Across Andhra Pradesh', font_size=18,
                bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.3),
                'iOS + Android  |  English + Telugu  |  Built for 45,000+ Schools',
                font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_footer(slide, 12, TOTAL_SLIDES)

    # Save
    prs.save(output_path)
    return output_path


def main():
    print("=" * 60)
    print("Vidya Nirmaan \u2014 Pitch Deck PPTX Generator")
    print("=" * 60)

    print("\n[1/6] Generating problem visual...")
    d1 = generate_problem_visual()

    print("[2/6] Generating solution overview...")
    d2 = generate_solution_overview()

    print("[3/6] Generating AI models visual...")
    d3 = generate_ai_models()

    print("[4/6] Generating pipeline visual...")
    d4 = generate_pipeline_visual()

    print("[5/6] Generating impact charts...")
    d5 = generate_impact_chart()

    print("[6/6] Generating tech stack visual...")
    d6 = generate_tech_stack_visual()

    diagrams = {
        'problem': d1, 'solution': d2, 'ai_models': d3,
        'pipeline': d4, 'impact': d5, 'tech_stack': d6,
    }

    print("\nAssembling PPTX...")
    output = build_pptx(diagrams)
    file_size = os.path.getsize(output) / (1024 * 1024)
    print(f"\n{'=' * 60}")
    print(f"Pitch Deck PPTX generated successfully!")
    print(f"  File: {output}")
    print(f"  Size: {file_size:.1f} MB")
    print(f"{'=' * 60}")

    import shutil
    shutil.rmtree(TEMP_DIR, ignore_errors=True)
    return output


if __name__ == '__main__':
    main()
